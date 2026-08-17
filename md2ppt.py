#!/usr/bin/env python3
"""Builds a PowerPoint deck from a markdown file using gartner_template.pptx's
real slide layouts (placeholders, theme, fonts) instead of hand-drawn shapes.
"""
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

TEMPLATE_PATH = Path(__file__).parent / "gartner_template.pptx"
DEFAULT_TABLE_STYLE_ID = "{6E25E649-3F16-4E02-A733-19D2CDBF48F0}"  # template's own default table style

# idx -> left-to-right column order used by the N-column layouts.
COLUMN_PLACEHOLDER_IDX = [1, 2, 10, 11]


# ---- markdown parsing (same grammar as the original md2ppt.js) ----------

def parse_meta(section):
    """Strips a leading 'key: value' block off a section. Returns (meta, remaining_body)."""
    meta = {}
    lines = section.strip().split("\n")
    while lines and re.match(r"^[\w-]+:", lines[0]):
        match = re.match(r"^([\w-]+):\s*(.*)$", lines.pop(0))
        meta[match.group(1)] = match.group(2)
    return meta, "\n".join(lines).strip()


def parse_deck(markdown):
    sections = [s.strip() for s in re.split(r"^---\s*$", markdown.strip(), flags=re.M) if s.strip()]
    deck_meta = {}
    if sections and "\ntype:" not in sections[0]:
        deck_meta, _ = parse_meta(sections.pop(0))
    slides = []
    for section in sections:
        meta, body = parse_meta(section)
        meta["body"] = body
        slides.append(meta)
    return deck_meta, slides


def parse_bullets(body):
    lines = (body or "").split("\n")
    return [re.sub(r"^\s*[-*]\s+", "", line) for line in lines if re.match(r"^\s*[-*]\s+", line)]


def parse_columns(body, max_columns):
    columns = [c.strip() for c in re.split(r"^::: column\s*$", body or "", flags=re.M) if c.strip()]
    return columns[:max_columns], len(columns)


def parse_table(body):
    rows = [line.strip() for line in (body or "").split("\n") if line.strip().startswith("|")]
    if len(rows) < 2:
        return None
    split_row = lambda line: [c.strip() for c in line.strip("|").split("|")]
    header = split_row(rows[0])
    data_rows = [split_row(r) for r in rows[2:]]
    return header, data_rows


# ---- pptx helpers ----------------------------------------------------------

def set_no_bullet(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def fill_bullets(placeholder, items, fallback_text=""):
    tf = placeholder.text_frame
    tf.clear()
    if not items:
        tf.paragraphs[0].text = fallback_text
        return
    tf.paragraphs[0].text = items[0]
    for item in items[1:]:
        tf.add_paragraph().text = item


def fill_column(placeholder, column_text):
    lines = column_text.split("\n")
    heading = re.sub(r"^#+\s*", "", lines[0]) if lines and lines[0].startswith("#") else None
    body = "\n".join(lines[1:]) if heading else column_text
    items = parse_bullets(body)

    tf = placeholder.text_frame
    tf.clear()
    first = True
    if heading:
        tf.paragraphs[0].text = heading
        tf.paragraphs[0].font.bold = True
        set_no_bullet(tf.paragraphs[0])
        first = False
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = item
        first = False


def add_table(slide, left, top, width, height, header, data_rows):
    n_rows = 1 + len(data_rows)
    n_cols = len(header)
    row_height = min(height // max(n_rows, 1), 500000)  # cap row height so tables don't stretch
    graphic_frame = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_height * n_rows)
    table = graphic_frame.table
    table._tbl.find(qn("a:tblPr")).set("firstRow", "1")
    table._tbl.find(qn("a:tblPr")).set("bandRow", "1")
    table._tbl.find(qn("a:tblPr")).find(qn("a:tableStyleId")).text = DEFAULT_TABLE_STYLE_ID
    for c, text in enumerate(header):
        table.cell(0, c).text = text
    for r, row in enumerate(data_rows, start=1):
        for c, text in enumerate(row):
            if c < n_cols:
                table.cell(r, c).text = text


# ---- slide builders ---------------------------------------------------------

def build_title(slide, layout, slide_data, deck_meta):
    slide.placeholders[0].text = slide_data.get("title") or deck_meta.get("title", "")
    slide.placeholders[1].text = slide_data.get("subtitle") or deck_meta.get("subtitle", "")


def build_bullets(slide, layout, slide_data, deck_meta):
    slide.placeholders[0].text = slide_data.get("title", "")
    items = parse_bullets(slide_data.get("body"))
    fill_bullets(slide.placeholders[1], items, slide_data.get("body", ""))


def build_columns(n):
    def _build(slide, layout, slide_data, deck_meta):
        slide.placeholders[0].text = slide_data.get("title", "")
        columns, total = parse_columns(slide_data.get("body"), n)
        if total > n:
            print(f"Warning: slide '{slide_data.get('title', '')}' has {total} '::: column' blocks; "
                  f"only the first {n} are rendered.", file=sys.stderr)
        for idx, column_text in zip(COLUMN_PLACEHOLDER_IDX[:n], columns):
            fill_column(slide.placeholders[idx], column_text)
    return _build


def build_takeaway(slide, layout, slide_data, deck_meta):
    # "Quote" layout: placeholder 0 is the large statement, placeholder 1 is a small attribution line.
    slide.placeholders[0].text = slide_data.get("body", "")
    slide.placeholders[1].text = slide_data.get("title", "")


def build_divider(slide, layout, slide_data, deck_meta):
    slide.placeholders[0].text = slide_data.get("title", "")


def build_table(slide, layout, slide_data, deck_meta):
    slide.placeholders[0].text = slide_data.get("title", "")
    parsed = parse_table(slide_data.get("body"))
    content_ph = slide.placeholders[1]
    if not parsed:
        print(f"Warning: slide '{slide_data.get('title', '')}' has type: table but no markdown table "
              f"found in body; rendering as plain text.", file=sys.stderr)
        content_ph.text_frame.text = slide_data.get("body", "")
        return
    header, data_rows = parsed
    left, top, width, height = content_ph.left, content_ph.top, content_ph.width, content_ph.height
    sp = content_ph._element
    sp.getparent().remove(sp)
    add_table(slide, left, top, width, height, header, data_rows)


LAYOUT_NAMES = {
    "title": "Title Slide",
    "bullets": "Title and Content",
    "two-column": "Two column",
    "three-column": "Three column",
    "four-column": "Four column",
    "takeaway": "Quote",
    "divider": "Divider Slide",
    "table": "Title and Content",
}

BUILDERS = {
    "title": build_title,
    "bullets": build_bullets,
    "two-column": build_columns(2),
    "three-column": build_columns(3),
    "four-column": build_columns(4),
    "takeaway": build_takeaway,
    "divider": build_divider,
    "table": build_table,
}


def build(input_path, output_path):
    try:
        markdown = Path(input_path).read_text(encoding="utf-8")
    except OSError as err:
        raise SystemExit(f"Could not read input file \"{input_path}\": {err}")

    deck_meta, slides = parse_deck(markdown)
    if not slides:
        raise SystemExit(f'No slides found in "{input_path}". '
                          f'Each slide must be separated by a line containing only "---".')

    prs = Presentation(str(TEMPLATE_PATH))
    prs.core_properties.author = deck_meta.get("presenter", "")
    prs.core_properties.title = deck_meta.get("title", "")

    layouts_by_name = {layout.name: layout for layout in prs.slide_masters[0].slide_layouts}
    # Start from a blank deck: drop the template's own instructional example slides
    # (not just the <p:sldId> reference -- also the underlying part, or saving
    # collides on the reused slideN.xml partnames).
    xml_slides = prs.slides._sldIdLst
    for sld_id in list(xml_slides):
        prs.part.drop_rel(sld_id.rId)
        xml_slides.remove(sld_id)

    for slide_data in slides:
        slide_type = slide_data.get("type") or "bullets"
        builder = BUILDERS.get(slide_type, build_bullets)
        layout_name = LAYOUT_NAMES.get(slide_type, "Title and Content")
        layout = layouts_by_name[layout_name]
        slide = prs.slides.add_slide(layout)
        builder(slide, layout, slide_data, deck_meta)

    prs.save(output_path)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python md2ppt.py input.md output.pptx", file=sys.stderr)
        sys.exit(2)
    build(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else "output.pptx")
